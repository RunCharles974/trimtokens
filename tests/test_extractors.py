"""Tests d'intégration Phase 2 : extracteurs texte natif + core dispatch.

Les fixtures sont générées programmatiquement (via python-docx, python-pptx,
openpyxl) pour rester self-contained et indépendants du système de fichiers.
"""

from __future__ import annotations

from pathlib import Path

import pytest

from trimtokens.models import ExtractOptions

# --- TXT / MD ----------------------------------------------------------------


def test_txt_extractor_basic(tmp_path: Path) -> None:
    from trimtokens.extractors.txt_md import extract

    f = tmp_path / "note.txt"
    f.write_text("Hello world.\nLigne 2.", encoding="utf-8")

    doc = extract(f, ExtractOptions())
    assert doc.source_type == "txt"
    assert doc.source_path == f
    assert len(doc.sections) == 1
    assert "Hello world." in doc.sections[0].content
    assert "Ligne 2." in doc.sections[0].content


def test_txt_extractor_handles_utf8_bom(tmp_path: Path) -> None:
    from trimtokens.extractors.txt_md import extract

    f = tmp_path / "bom.txt"
    f.write_bytes(b"\xef\xbb\xbfHello with BOM")

    doc = extract(f, ExtractOptions())
    assert "Hello with BOM" in doc.sections[0].content
    assert doc.metadata["encoding"] == "utf-8-sig"


def test_md_extractor_marks_source_type_markdown(tmp_path: Path) -> None:
    from trimtokens.extractors.txt_md import extract

    f = tmp_path / "doc.md"
    f.write_text("# Titre\n\nCorps.", encoding="utf-8")

    doc = extract(f, ExtractOptions())
    assert doc.source_type == "markdown"


def test_txt_extractor_falls_back_to_chardet_on_latin1(tmp_path: Path) -> None:
    from trimtokens.extractors.txt_md import extract

    f = tmp_path / "latin.txt"
    f.write_bytes("Café déjà vu".encode("latin-1"))

    doc = extract(f, ExtractOptions())
    assert "Caf" in doc.sections[0].content


# --- HTML --------------------------------------------------------------------


def test_html_extractor_strips_scripts_and_styles(tmp_path: Path) -> None:
    from trimtokens.extractors.html import extract

    html = """<html>
    <head><title>Mon Titre</title><style>body{color:red}</style></head>
    <body>
      <h1>Heading principal</h1>
      <p>Un paragraphe utile.</p>
      <script>evilCode();</script>
    </body>
    </html>"""
    f = tmp_path / "page.html"
    f.write_text(html, encoding="utf-8")

    doc = extract(f, ExtractOptions())
    assert doc.source_type == "html"
    assert doc.title == "Mon Titre"
    content = doc.sections[0].content
    assert "Heading principal" in content
    assert "Un paragraphe utile." in content
    assert "evilCode" not in content
    assert "color:red" not in content


def test_html_extractor_no_title(tmp_path: Path) -> None:
    from trimtokens.extractors.html import extract

    f = tmp_path / "notitle.html"
    f.write_text("<html><body><p>Juste un paragraphe.</p></body></html>", encoding="utf-8")

    doc = extract(f, ExtractOptions())
    assert doc.title is None
    assert "Juste un paragraphe." in doc.sections[0].content


# --- DOCX --------------------------------------------------------------------


def test_docx_extractor_preserves_headings_and_paragraphs(tmp_path: Path) -> None:
    from docx import Document  # type: ignore[import-untyped]

    from trimtokens.extractors.docx import extract

    docx_path = tmp_path / "report.docx"
    doc_in = Document()
    doc_in.add_heading("Titre principal", level=1)
    doc_in.add_heading("Section A", level=2)
    doc_in.add_paragraph("Premier paragraphe.")
    doc_in.add_heading("Sous-section", level=3)
    doc_in.add_paragraph("Texte intermédiaire.")
    doc_in.save(str(docx_path))

    doc = extract(docx_path, ExtractOptions())
    assert doc.source_type == "docx"
    content = doc.sections[0].content
    assert "# Titre principal" in content
    assert "## Section A" in content
    assert "### Sous-section" in content
    assert "Premier paragraphe." in content
    assert "Texte intermédiaire." in content


def test_docx_extractor_renders_tables(tmp_path: Path) -> None:
    from docx import Document  # type: ignore[import-untyped]

    from trimtokens.extractors.docx import extract

    docx_path = tmp_path / "table.docx"
    doc_in = Document()
    doc_in.add_paragraph("Intro.")
    table = doc_in.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "Nom"
    table.cell(0, 1).text = "Âge"
    table.cell(1, 0).text = "Alice"
    table.cell(1, 1).text = "30"
    doc_in.save(str(docx_path))

    doc = extract(docx_path, ExtractOptions())
    content = doc.sections[0].content
    assert "| Nom | Âge |" in content
    assert "| --- | --- |" in content
    assert "| Alice | 30 |" in content


# --- PPTX --------------------------------------------------------------------


def test_pptx_extractor_one_section_per_slide(tmp_path: Path) -> None:
    from pptx import Presentation  # type: ignore[import-untyped]

    from trimtokens.extractors.pptx import extract

    pptx_path = tmp_path / "deck.pptx"
    prs = Presentation()
    layout = prs.slide_layouts[1]  # Title and Content

    slide1 = prs.slides.add_slide(layout)
    slide1.shapes.title.text = "Premier Slide"
    slide1.placeholders[1].text = "Point A\nPoint B"

    slide2 = prs.slides.add_slide(layout)
    slide2.shapes.title.text = "Deuxième Slide"
    slide2.notes_slide.notes_text_frame.text = "Note du présentateur."

    prs.save(str(pptx_path))

    doc = extract(pptx_path, ExtractOptions())
    assert doc.source_type == "pptx"
    assert len(doc.sections) == 2
    assert "Slide 1" in doc.sections[0].header
    assert "Premier Slide" in doc.sections[0].header
    assert "Point A" in doc.sections[0].content
    assert "### Notes du présentateur" in doc.sections[1].content
    assert "Note du présentateur." in doc.sections[1].content


# --- XLSX / CSV --------------------------------------------------------------


def test_xlsx_extractor_one_section_per_sheet(tmp_path: Path) -> None:
    from openpyxl import Workbook

    from trimtokens.extractors.xlsx_csv import extract

    xlsx_path = tmp_path / "data.xlsx"
    wb = Workbook()
    ws1 = wb.active
    assert ws1 is not None
    ws1.title = "Personnes"
    ws1.append(["Nom", "Âge"])
    ws1.append(["Alice", 30])
    ws1.append(["Bob", 25])

    ws2 = wb.create_sheet("Pays")
    ws2.append(["Code", "Nom"])
    ws2.append(["FR", "France"])
    wb.save(str(xlsx_path))

    doc = extract(xlsx_path, ExtractOptions())
    assert doc.source_type == "xlsx"
    assert len(doc.sections) == 2

    headers = [s.header for s in doc.sections]
    assert "Feuille : Personnes" in headers
    assert "Feuille : Pays" in headers

    personnes = next(s for s in doc.sections if s.header == "Feuille : Personnes")
    assert "| Nom | Âge |" in personnes.content
    assert "| Alice | 30 |" in personnes.content


def test_xlsx_extractor_skips_empty_sheets(tmp_path: Path) -> None:
    from openpyxl import Workbook

    from trimtokens.extractors.xlsx_csv import extract

    xlsx_path = tmp_path / "with_empty.xlsx"
    wb = Workbook()
    ws_empty = wb.active
    assert ws_empty is not None
    ws_empty.title = "Vide"
    ws_full = wb.create_sheet("Avec données")
    ws_full.append(["A", "B"])
    ws_full.append([1, 2])
    wb.save(str(xlsx_path))

    doc = extract(xlsx_path, ExtractOptions())
    assert len(doc.sections) == 1
    assert doc.sections[0].header == "Feuille : Avec données"


def test_csv_extractor(tmp_path: Path) -> None:
    from trimtokens.extractors.xlsx_csv import extract

    csv_path = tmp_path / "data.csv"
    csv_path.write_text("Nom,Âge\nAlice,30\nBob,25\n", encoding="utf-8")

    doc = extract(csv_path, ExtractOptions())
    assert doc.source_type == "csv"
    assert len(doc.sections) == 1
    assert doc.sections[0].header == "Feuille : data"
    assert "Alice" in doc.sections[0].content
    assert "| --- | --- |" in doc.sections[0].content


# --- RTF ---------------------------------------------------------------------


def test_rtf_extractor(tmp_path: Path) -> None:
    from trimtokens.extractors.rtf import extract

    rtf = r"{\rtf1\ansi\deff0 {\fonttbl {\f0 Times;}} Hello\par This is RTF content.}"
    f = tmp_path / "doc.rtf"
    f.write_text(rtf, encoding="utf-8")

    doc = extract(f, ExtractOptions())
    assert doc.source_type == "rtf"
    content = doc.sections[0].content
    assert "Hello" in content
    assert "RTF content" in content


# --- Core dispatch -----------------------------------------------------------


def test_core_process_txt_returns_markdown_with_frontmatter(tmp_path: Path) -> None:
    from trimtokens.core import process

    f = tmp_path / "note.txt"
    f.write_text("Bonjour\n\nMonde.", encoding="utf-8")

    result = process(f)
    assert result.markdown.startswith("---\n")
    assert "source: note.txt" in result.markdown
    assert "type: txt" in result.markdown
    assert "Bonjour" in result.markdown
    assert result.stats.original_size_bytes > 0
    assert result.stats.cleaned_chars > 0


def test_core_process_csv_dispatches_correctly(tmp_path: Path) -> None:
    from trimtokens.core import process

    f = tmp_path / "data.csv"
    f.write_text("a,b\n1,2", encoding="utf-8")

    result = process(f)
    assert result.document.source_type == "csv"
    assert "| a | b |" in result.markdown


def test_core_process_unknown_extension_raises(tmp_path: Path) -> None:
    from trimtokens.core import process

    f = tmp_path / "weird.xyz"
    f.write_text("data", encoding="utf-8")

    with pytest.raises(ValueError, match="Format non supporté"):
        process(f)


def test_core_process_missing_file_raises(tmp_path: Path) -> None:
    from trimtokens.core import process

    with pytest.raises(FileNotFoundError):
        process(tmp_path / "does_not_exist.txt")


def test_core_process_cleans_extra_whitespace(tmp_path: Path) -> None:
    from trimtokens.core import process

    f = tmp_path / "messy.txt"
    f.write_text("hello    world\n\n\n\n\nfoo", encoding="utf-8")

    result = process(f)
    # Pipeline doit avoir collapsé espaces et sauts de ligne
    assert "    " not in result.document.sections[0].content
    assert "\n\n\n" not in result.document.sections[0].content
