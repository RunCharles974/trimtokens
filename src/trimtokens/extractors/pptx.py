"""Extracteur PPTX via python-pptx.

Pour chaque slide :
- Header `Slide N — Titre` (titre extrait du placeholder titre)
- Contenu des autres shapes textuels concaténés
- Bloc `### Notes du présentateur` si notes présentes
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation  # type: ignore[import-untyped]

from trimtokens.models import ExtractedDocument, ExtractOptions, Section


def extract(path: Path, options: ExtractOptions) -> ExtractedDocument:
    presentation = Presentation(str(path))

    doc_title = presentation.core_properties.title or None

    sections: list[Section] = []
    for idx, slide in enumerate(presentation.slides, start=1):
        slide_title, body_text = _extract_slide_text(slide)
        notes_text = _extract_notes(slide)

        header = f"Slide {idx}"
        if slide_title:
            header += f" — {slide_title}"

        content_parts: list[str] = []
        if body_text:
            content_parts.append(body_text)
        if notes_text:
            content_parts.append(f"### Notes du présentateur\n\n{notes_text}")

        sections.append(Section(header=header, content="\n\n".join(content_parts)))

    return ExtractedDocument(
        source_path=path,
        source_type="pptx",
        title=doc_title,
        sections=sections,
    )


def _extract_slide_text(slide: object) -> tuple[str, str]:
    """Retourne (titre_slide, corps_concat)."""
    title = ""
    title_shape = getattr(slide.shapes, "title", None)  # type: ignore[attr-defined]
    if title_shape is not None and title_shape.has_text_frame:
        title = title_shape.text_frame.text.strip()

    body_parts: list[str] = []
    for shape in slide.shapes:  # type: ignore[attr-defined]
        if not shape.has_text_frame:
            continue
        if shape is title_shape:
            continue
        text = shape.text_frame.text.strip()
        if text:
            body_parts.append(text)

    return title, "\n\n".join(body_parts)


def _extract_notes(slide: object) -> str:
    if not getattr(slide, "has_notes_slide", False):
        return ""
    notes_slide = slide.notes_slide  # type: ignore[attr-defined]
    if notes_slide is None:
        return ""
    notes_tf = notes_slide.notes_text_frame
    if notes_tf is None:
        return ""
    return notes_tf.text.strip()
